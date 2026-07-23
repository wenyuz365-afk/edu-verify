"""仿手绘流程图 -- 传统流程图风格"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

BG = RGBColor(0x0F,0x17,0x2A)
WHITE = RGBColor(0xF1,0xF5,0xF9); GRAY = RGBColor(0x94,0xA3,0xB8)
GREEN = RGBColor(0x22,0xC5,0x5E); ORANGE = RGBColor(0xF5,0x9E,0x0B)
RED = RGBColor(0xEF,0x44,0x44); BLUE = RGBColor(0x3B,0x82,0xF6)
CYAN = RGBColor(0x06,0xB6,0xD4); PURPLE = RGBColor(0xA8,0x55,0xF7)

def set_bg(s, c): s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def rect(s, l, t, w, h, fill=None, border=None, text="", color=WHITE, size=11, bold=False):
    fill = fill or RGBColor(0x1E,0x29,0x3B)
    border = border or RGBColor(0x47,0x55,0x69)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(2)
    if text:
        tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
        p.font.color.rgb = color; p.font.bold = bold; p.alignment = PP_ALIGN.CENTER
    return sh

def diamond(s, l, t, w, h, fill=None, border=None, text="", color=WHITE, size=10):
    fill = fill or RGBColor(0x1E,0x29,0x3B)
    border = border or RGBColor(0x47,0x55,0x69)
    sh = s.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(2)
    if text:
        tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
        p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    return sh

def txt(s, l, t, w, h, text, color=WHITE, size=11, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tb

def arrow_line(s, x1, y1, x2, y2, color=GRAY, w=2):
    conn = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color; conn.line.width = Pt(w)
    # add arrowhead
    conn.line._ln.append(
        qn('a:tailEnd')  # or headEnd
    )
    return conn

def vert_arrow(s, x, y1, y2, color=GRAY):
    """垂直箭头: x中心, y1顶部, y2底部"""
    sh = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x, y1, x, y2)
    sh.line.color.rgb = color; sh.line.width = Pt(2)
    return sh

def label(s, l, t, text, color=GRAY, size=9):
    txt(s, l, t, Inches(1.5), Inches(0.25), text, color, size, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# Main Flowchart Slide
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide, BG)
txt(slide, Inches(0), Inches(0.06), Inches(16), Inches(0.4), "学历信息 AI 核验系统 -- 完整业务流程图", WHITE, 24, True, PP_ALIGN.CENTER)

# ─── ROW 1: 上传层 ───
y1 = Inches(0.55)
rect(slide, Inches(2.0), y1, Inches(2.5), Inches(0.55), text="上传学位证", color=CYAN, size=14, bold=True)
rect(slide, Inches(5.0), y1, Inches(2.5), Inches(0.55), text="上传毕业证", color=CYAN, size=14, bold=True)
rect(slide, Inches(11.5), y1, Inches(2.5), Inches(0.55), text="上传简历", color=CYAN, size=14, bold=True)

# Arrows from upload to OCR
vert_arrow(slide, Inches(3.25), y1+Inches(0.55), y1+Inches(0.95))
vert_arrow(slide, Inches(6.25), y1+Inches(0.55), y1+Inches(0.95))
vert_arrow(slide, Inches(12.75), y1+Inches(0.55), y1+Inches(1.3))

# ─── ROW 2: OCR + LLM ───
y2 = Inches(1.0)
rect(slide, Inches(1.5), y2, Inches(5.0), Inches(0.65), fill=RGBColor(0x0E,0x3A,0x5C), border=CYAN, text="PaddleOCR 图片文字识别", color=CYAN, size=15, bold=True)
# Sub-descriptions
txt(slide, Inches(1.0), y2+Inches(0.7), Inches(6), Inches(0.25), "去印章 | 自适应缩放 | CLAHE增强 | 双通道", GRAY, 9, False, PP_ALIGN.CENTER)

# Arrow from OCR to LLM
vert_arrow(slide, Inches(4.0), y2+Inches(0.65), y2+Inches(1.15))

y3 = Inches(1.8)
rect(slide, Inches(1.5), y3, Inches(5.0), Inches(0.65), fill=RGBColor(0x0E,0x3A,0x5C), border=CYAN, text="DeepSeek V4 Pro 字段提取", color=CYAN, size=15, bold=True)
txt(slide, Inches(1.0), y3+Inches(0.7), Inches(6), Inches(0.25), "姓名/学校/专业/学历层次/毕业日期/证书编号/学习形式", GRAY, 9, False, PP_ALIGN.CENTER)

# Also resume OCR -> LLM
rect(slide, Inches(10.5), y2, Inches(4.5), Inches(0.55), fill=RGBColor(0x0E,0x3A,0x5C), border=CYAN, text="简历OCR识别 + LLM提取", color=CYAN, size=14, bold=True)
txt(slide, Inches(10.5), y2+Inches(0.6), Inches(4.5), Inches(0.3), "多段教育经历: 学校/专业/学历/起止年份/状态", GRAY, 9, False, PP_ALIGN.CENTER)

# Vertical line from resume -> merge
vert_arrow(slide, Inches(12.75), y2+Inches(0.55), y3)

# ─── Connection arrows ───
# From LLM output to verification area
vert_arrow(slide, Inches(4.0), y3+Inches(0.65), y3+Inches(1.05))

# ─── ROW 3: 提取结果输出 ───
y4 = Inches(3.15)
rect(slide, Inches(1.5), y4, Inches(11.0), Inches(0.5), fill=RGBColor(0x1A,0x3A,0x2A), border=GREEN, text="结构化字段输出: 姓名 | 性别 | 出生日期 | 学校 | 专业 | 学历层次 | 学位 | 毕业日期 | 证书编号", color=GREEN, size=11, bold=True)

vert_arrow(slide, Inches(7.0), y4+Inches(0.5), y4+Inches(0.9))

# ─── ROW 4: 核验检查项 (菱形+矩形组合) ───
y5 = Inches(4.15)

# Title
txt(slide, Inches(0.3), y5, Inches(2), Inches(0.35), "核验引擎:", WHITE, 13, True)

# Check items in a flow
checks = [
    ("教育部名单\n院校真实性\n(学校是否在名单中)", RED),
    ("学信库查询\n学籍信息比对\n(查该学生记录)", RED),
    ("同编号重复\n检测", ORANGE),
    ("毕业年龄\n异常检测", ORANGE),
    ("院校-专业\n一致性", ORANGE),
    ("ELA图片\n篡改检测", RED),
    ("黑名单\n筛查", RED),
]

cx = Inches(2.2)
for i, (name, color) in enumerate(checks):
    diamond(slide, cx, y5, Inches(1.6), Inches(1.0), border=color, text=name, color=color, size=8)
    cx += Inches(1.75)

# Arrows between checks
for i in range(len(checks)-1):
    ax = Inches(2.2) + (i+1)*Inches(1.75) - Inches(0.15)
    txt(slide, ax, y5+Inches(0.35), Inches(0.3), Inches(0.3), ">", GRAY, 14, False, PP_ALIGN.CENTER)

# ─── ROW 5: 交叉验证 ───
y6 = Inches(5.4)
rect(slide, Inches(1.5), y6, Inches(5.0), Inches(0.7), border=ORANGE, text="双证交叉比对\n学位证 vs 毕业证: 姓名/学校/专业/学历/日期", color=ORANGE, size=11, bold=True)
rect(slide, Inches(7.0), y6, Inches(5.5), Inches(0.7), border=ORANGE, text="简历交叉验证\n按学历层级逐段比对(学校+专业包含匹配+年份严格)", color=ORANGE, size=11, bold=True)

# ─── ROW 6: 判定菱形 ───
y7 = Inches(6.5)
diamond(slide, Inches(3.0), y7, Inches(2.5), Inches(1.2), border=ORANGE, text="系统综合判定\nPASS / REVIEW / ALERT", color=ORANGE, size=10)

# Connected arrows to diamond
vert_arrow(slide, Inches(3.5), y6+Inches(0.7), y7)

# ─── ROW 7: 三分流结果 ───
y8 = Inches(7.95)

# Left: PASS
rect(slide, Inches(0.5), y8, Inches(3.5), Inches(0.7), fill=RGBColor(0x14,0x4E,0x2A), border=GREEN, text="PASS 通过", color=GREEN, size=18, bold=True)
txt(slide, Inches(0.5), y8+Inches(0.75), Inches(3.5), Inches(0.3), "全部检查无异常 -> 自动进入通过名单", GREEN, 10, False, PP_ALIGN.CENTER)

# Center: REVIEW
rect(slide, Inches(4.5), y8, Inches(3.5), Inches(0.7), fill=RGBColor(0x5C,0x3A,0x0E), border=ORANGE, text="REVIEW 复核", color=ORANGE, size=18, bold=True)
txt(slide, Inches(4.5), y8+Inches(0.75), Inches(3.5), Inches(0.3), "有可疑项 -> 进入人工核验台", ORANGE, 10, False, PP_ALIGN.CENTER)

# Right: ALERT
rect(slide, Inches(8.5), y8, Inches(3.5), Inches(0.7), fill=RGBColor(0x5C,0x14,0x14), border=RED, text="ALERT 异常预警", color=RED, size=18, bold=True)
txt(slide, Inches(8.5), y8+Inches(0.75), Inches(3.5), Inches(0.3), "严重异常 -> 进入人工核验台，优先处理", RED, 10, False, PP_ALIGN.CENTER)

# Arrows from diamond to three outcomes
vert_arrow(slide, Inches(1.5), y7+Inches(1.2), y8, GREEN)
vert_arrow(slide, Inches(5.0), y7+Inches(1.2), y8, ORANGE)
vert_arrow(slide, Inches(9.0), y7+Inches(1.2), y8, RED)

# ─── ROW 8: HR 判定 ───
y9 = Inches(9.7)  # off slide, but close

# HR判定 area on right side
txt(slide, Inches(11.5), y4+Inches(0.5), Inches(4), Inches(0.3), "人工核验台 (HR判定)", WHITE, 14, True, PP_ALIGN.CENTER)
rect(slide, Inches(11.5), y4+Inches(0.9), Inches(4.0), Inches(2.8), border=WHITE)

txt(slide, Inches(11.7), y4+Inches(1.0), Inches(3.6), Inches(0.8), "HR查看原始证书图片\n及系统判定原因\n做出最终决定:", GRAY, 10)

# Three HR buttons
rect(slide, Inches(11.8), y4+Inches(1.9), Inches(3.3), Inches(0.55), fill=GREEN, border=GREEN, text="PASS 通过", color=WHITE, size=13, bold=True)
rect(slide, Inches(11.8), y4+Inches(2.55), Inches(3.3), Inches(0.55), fill=RED, border=RED, text="FAIL 淘汰", color=WHITE, size=13, bold=True)
rect(slide, Inches(11.8), y4+Inches(3.2), Inches(3.3), Inches(0.45), fill=ORANGE, border=ORANGE, text="HOLD 待补充材料", color=WHITE, size=12, bold=True)

# ─── 数据看板 Footer ───
y10 = Inches(9.0)
rect(slide, Inches(1.5), y10, Inches(9.5), Inches(0.5), fill=RGBColor(0x1E,0x29,0x3B), border=BLUE, text="数据看板: 全部 | 通过 | 不通过 | 待定 | 待核验  |  内联修改状态  |  CSV导出", color=BLUE, size=11, bold=True)

# Legend at bottom
txt(slide, Inches(0.3), Inches(0), Inches(0), Inches(0), "", GRAY, 1)  # dummy

prs.save("D:/edu-verify/学历核验系统流程图_仿手绘.pptx")
print("Saved: 仿手绘流程图")
