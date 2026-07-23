"""学历核验系统PPT -- 演示版，大框架醒目，操作选项清晰"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors -- 柔和不刺眼
BG = RGBColor(0x1A, 0x1F, 0x2E)
WHITE = RGBColor(0xE2, 0xE8, 0xF0); GRAY = RGBColor(0x7A, 0x86, 0x99)
GREEN = RGBColor(0x4A, 0x9E, 0x6E); ORANGE = RGBColor(0xC8, 0x8A, 0x42)
RED = RGBColor(0xC0, 0x55, 0x55); BLUE = RGBColor(0x5D, 0x8C, 0xC4)
CYAN = RGBColor(0x4A, 0x9B, 0xA8); PURPLE = RGBColor(0x8A, 0x6D, 0xC0)
PINK = RGBColor(0xC0, 0x6A, 0x82)

def set_slide_bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, border=None):
    fill = fill or RGBColor(0x24,0x2B,0x3D)
    border = border or RGBColor(0x33,0x41,0x55)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(1.5)
    return s

def big_header(slide, l, t, text, color, icon=""):
    """大框架标题 -- 醒目"""
    s = box(slide, l, t, Inches(14), Inches(0.65), color, color)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"  {icon} {text}" if icon else f"    {text}"
    p.font.size = Pt(18); p.font.color.rgb = WHITE; p.font.bold = True

def txt(slide, l, t, w, h, text, color=WHITE, size=12, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tb

def option_tag(slide, l, t, text, color):
    """选项标签"""
    s = box(slide, l, t, Inches(2.2), Inches(0.4), color, color)
    tf = s.text_frame; p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

def check_card(slide, l, t, text, color=WHITE, w=Inches(4.3)):
    s = box(slide, l, t, w, Inches(0.55))
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(10); p.font.color.rgb = color
    return s

def arrow(slide, l, t):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, Inches(0.22), Inches(0.3))
    s.fill.solid(); s.fill.fore_color.rgb = GRAY; s.line.fill.background()

def hr_pill(slide, l, t, text, color):
    s = box(slide, l, t, Inches(2.8), Inches(0.65), color, color)
    tf = s.text_frame; p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(22); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

# ═══════════════════ SLIDE 1: 上传层 + 识别层 ═══════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s1, BG)
txt(s1, Inches(0), Inches(0.08), Inches(16), Inches(0.5), "学历信息 AI 核验系统 -- 操作流程", WHITE, 26, True, PP_ALIGN.CENTER)
txt(s1, Inches(0), Inches(0.5), Inches(16), Inches(0.3), "面向HR用户 | 三步完成证书真伪核验", GRAY, 12, False, PP_ALIGN.CENTER)

# === 第一步: 选择招聘场景 ===
y = Inches(1.0)
big_header(s1, Inches(1), y, "第一步: 选择招聘场景", BLUE)
y += Inches(0.8)
box(s1, Inches(1), y, Inches(14), Inches(1.0))
txt(s1, Inches(1.3), y+Inches(0.08), Inches(13), Inches(0.35), "根据招聘类型，系统自动切换核验规则", WHITE, 14, True)
option_tag(s1, Inches(1.5), y+Inches(0.5), "校  招", CYAN)
txt(s1, Inches(4.0), y+Inches(0.5), Inches(3), Inches(0.35), "应届毕业生\n有证验真伪，无证标「待取证」", GRAY, 10)
option_tag(s1, Inches(5.8), y+Inches(0.5), "社  招", BLUE)
txt(s1, Inches(8.3), y+Inches(0.5), Inches(3), Inches(0.35), "已毕业候选人\n所有学历必须有证书", GRAY, 10)
option_tag(s1, Inches(10.1), y+Inches(0.5), "实  习", PINK)
txt(s1, Inches(12.6), y+Inches(0.5), Inches(2.2), Inches(0.35), "在读学生\n已毕业学历必验证", GRAY, 10)
arrow(s1, Inches(7.9), y+Inches(1.05))
y += Inches(1.25)

# === 第二步: 上传材料 ===
big_header(s1, Inches(1), y, "第二步: 上传材料", BLUE)
y += Inches(0.8)
box(s1, Inches(1), y, Inches(14), Inches(2.0))
txt(s1, Inches(1.3), y+Inches(0.05), Inches(13), Inches(0.3), "上传模式", WHITE, 14, True)
option_tag(s1, Inches(1.5), y+Inches(0.45), "单人上传", CYAN)
txt(s1, Inches(4.0), y+Inches(0.45), Inches(10.5), Inches(0.45), "动态添加学历层级(本科/硕士/博士)，每层级上传学位证+毕业证(毕业证可选)  |  简历为必传项", GRAY, 11)
option_tag(s1, Inches(1.5), y+Inches(0.95), "批量上传", CYAN)
txt(s1, Inches(4.0), y+Inches(0.95), Inches(10.5), Inches(0.45), "上传ZIP压缩包，每个子文件夹对应一位候选人，系统自动按文件名关键词分类学历层次", GRAY, 11)

txt(s1, Inches(1.3), y+Inches(1.5), Inches(13), Inches(0.3), "支持文件格式: JPG / PNG / PDF   |   PDF自动转换为图片   |   证书原件自动存档供人工核验查看", GRAY, 10)
arrow(s1, Inches(7.9), y+Inches(2.05))
y += Inches(2.25)

# === 第三步: 开始核验 ===
big_header(s1, Inches(1), y, "第三步: 开始核验 (系统自动完成)", BLUE)
y += Inches(0.8)
box(s1, Inches(1), y, Inches(14), Inches(1.3))

# OCR box
box(s1, Inches(1.3), y+Inches(0.1), Inches(6.5), Inches(1.1), border=CYAN)
txt(s1, Inches(1.5), y+Inches(0.15), Inches(6), Inches(0.25), "OCR 文字识别", CYAN, 14, True)
txt(s1, Inches(1.5), y+Inches(0.45), Inches(6), Inches(0.6), "PaddleOCR 2.6.2 本地引擎\n自动去除红色印章干扰 -> 自适应缩放清晰度增强 -> 输出全文文本\n证书和简历的图片文字全部提取", GRAY, 10)

# LLM box
box(s1, Inches(8.2), y+Inches(0.1), Inches(6.5), Inches(1.1), border=CYAN)
txt(s1, Inches(8.4), y+Inches(0.15), Inches(6), Inches(0.25), "AI 字段提取", CYAN, 14, True)
txt(s1, Inches(8.4), y+Inches(0.45), Inches(6), Inches(0.6), "DeepSeek V4 Pro 大语言模型\n理解证书文字 -> 提取: 姓名/学校/专业/学历/日期/编号等\n理解简历文字 -> 提取: 多段教育经历 + 在读/毕业状态", GRAY, 10)

arrow(s1, Inches(7.9), y+Inches(1.35))
y += Inches(1.6)

# Footer
txt(s1, Inches(1), y+Inches(0.1), Inches(14), Inches(0.3), "核验结果展示在三步流的结果页: 总览表 + 逐层详情 + 简历交叉验证   |   核验下一个候选人 一键返回上传页", GRAY, 11)

# ═══════════════════ SLIDE 2: 核验层 ═══════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s2, BG)
txt(s2, Inches(0), Inches(0.08), Inches(16), Inches(0.5), "学历信息 AI 核验系统 -- 核验引擎", WHITE, 26, True, PP_ALIGN.CENTER)
txt(s2, Inches(0), Inches(0.5), Inches(16), Inches(0.3), "系统自动完成的全部检查项 | 单证核验 + 交叉验证", GRAY, 12, False, PP_ALIGN.CENTER)

# ─── 模块A: 单证核验 ───
y = Inches(1.0)
big_header(s2, Inches(1), y, "模块A: 单证单项核验 (每份证书独立进行)", PURPLE)
y += Inches(0.8)

checks_a = [
    ("教育部名单校验", "学校是否在583所国内本科+191所海外认证名单中", RED),
    ("学信库查询", "按身份证号/证书编号/姓名查询学籍记录并逐字段比对", RED),
    ("证书状态检测", "识别结业证书、肄业证书(非正常毕业)", ORANGE),
    ("毕业年龄检测", "本科<20或>30岁标记异常，硕博采用不同年龄范围", ORANGE),
    ("院校-专业一致性", "正向查找学校是否真实开设该专业(12,697条数据)", ORANGE),
    ("ELA图片篡改检测", "JPEG压缩误差分析: 检测图片是否经过PS修改", RED),
]
for i, (title, desc, color) in enumerate(checks_a):
    col = i % 3; row = i // 3
    cx = Inches(1) + col * Inches(4.6); cy = y + row * Inches(0.95)
    s = box(s2, cx, cy, Inches(4.2), Inches(0.8), border=color)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.12)
    p1 = tf.paragraphs[0]; p1.text = title; p1.font.size = Pt(11); p1.font.color.rgb = color; p1.font.bold = True
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(9); p2.font.color.rgb = GRAY

y += Inches(2.1)

# ─── 模块B: 交叉验证 ───
big_header(s2, Inches(1), y, "模块B: 多材料交叉验证 (跨文件信息对照)", PURPLE)
y += Inches(0.8)

# 双证比对
box(s2, Inches(1), y, Inches(6.5), Inches(1.15), border=ORANGE)
txt(s2, Inches(1.2), y+Inches(0.08), Inches(6), Inches(0.25), "双证交叉比对", ORANGE, 14, True)
txt(s2, Inches(1.2), y+Inches(0.4), Inches(6), Inches(0.7), "学位证书 vs 毕业证书，逐字段比对:\n姓名 / 学校 / 专业 / 学历层次 / 毕业日期\n任一项不一致 -> 标记异常", GRAY, 10)

# 简历交叉
box(s2, Inches(8.1), y, Inches(6.9), Inches(1.15), border=ORANGE)
txt(s2, Inches(8.3), y+Inches(0.08), Inches(6.5), Inches(0.25), "简历分层交叉验证", ORANGE, 14, True)
txt(s2, Inches(8.3), y+Inches(0.4), Inches(6.5), Inches(0.7), "简历提取多段教育经历 -> 按学历层级与对应证书逐段比对\n学校(简称/曾用名归一化) + 专业(包含匹配) + 年份(严格匹配)\n在读/退学/交换/辅修 -> 自动跳过，不要求证书", GRAY, 10)

# 特殊状态标签
option_tag(s2, Inches(1), y+Inches(1.3), "ongoing 跳过", CYAN)
option_tag(s2, Inches(3.5), y+Inches(1.3), "dropout 跳过", CYAN)
option_tag(s2, Inches(6.0), y+Inches(1.3), "exchange 跳过", CYAN)
option_tag(s2, Inches(8.5), y+Inches(1.3), "minor 跳过", CYAN)
option_tag(s2, Inches(11.0), y+Inches(1.3), "merged 正常验", CYAN)

y += Inches(1.85)
big_header(s2, Inches(1), y, "招聘类型判定差异 (缺证处理)", ORANGE)
y += Inches(0.8)
box(s2, Inches(1), y, Inches(14), Inches(0.85))
txt(s2, Inches(1.3), y+Inches(0.1), Inches(4), Inches(0.65), "校招: 有证验真伪 | 缺证=待取证(REVIEW) | 在读=跳过", CYAN, 12, True)
txt(s2, Inches(5.8), y+Inches(0.1), Inches(4), Inches(0.65), "社招: 有证验真伪 | 缺证=不通过(ALERT) | 在读=不存在", BLUE, 12, True)
txt(s2, Inches(10.3), y+Inches(0.1), Inches(4), Inches(0.65), "实习: 有证验真伪 | 缺证=不通过(ALERT) | 在读=跳过", PINK, 12, True)

# ═══════════════════ SLIDE 3: 判定 + 看板 ═══════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s3, BG)
txt(s3, Inches(0), Inches(0.08), Inches(16), Inches(0.5), "学历信息 AI 核验系统 -- 判定分流 + 人工核验", WHITE, 26, True, PP_ALIGN.CENTER)

y = Inches(1.0)
big_header(s3, Inches(1), y, "系统自动判定: 三分流", GREEN)
y += Inches(0.85)

# Three columns
for i, (label, color, desc, next_step) in enumerate([
    ("PASS 通过", GREEN,
     "所有检查均无异常\n教育部名单通过\n学信库查到且一致\n双证比对一致\n简历交叉无冲突\nELA无篡改痕迹\n年龄/专业正常",
     "自动进入通过名单\n无需HR介入"),
    ("REVIEW 需复核", ORANGE,
     "存在可疑项需要确认\n学信库字段不一致\n毕业年龄偏大/小\n专业未在列表中\nELA轻微异常\n双证/简历不一致\n校招待取证",
     "进入人工核验台\nHR查看材料后判定"),
    ("ALERT 告警", RED,
     "严重异常，优先处理\n学校不在教育部名单\n学信库查无此人\n社招缺证\n实习已毕业缺证\nELA明显篡改\n证书状态异常",
     "进入人工核验台\nHR重点核查"),
]):
    cx = Inches(1) + i * Inches(4.8)
    box(s3, cx, y, Inches(4.4), Inches(3.2), border=color)
    s = box(s3, cx, y, Inches(4.4), Inches(0.55), color, color)
    tf = s.text_frame; p = tf.paragraphs[0]
    p.text = label; p.font.size = Pt(18); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    txt(s3, cx+Inches(0.2), y+Inches(0.7), Inches(4), Inches(1.6), desc, GRAY, 10)
    # Arrow + next step
    txt(s3, cx+Inches(0.2), y+Inches(2.5), Inches(4), Inches(0.5), next_step, color, 12, True)

arrow(s3, Inches(3.2), y+Inches(3.25))
arrow(s3, Inches(8.0), y+Inches(3.25))
arrow(s3, Inches(12.8), y+Inches(3.25))
y += Inches(3.5)

# HR Review Station
big_header(s3, Inches(1), y, "人工核验台: HR 最终判定", ORANGE)
y += Inches(0.8)
box(s3, Inches(1), y, Inches(14), Inches(1.25))

txt(s3, Inches(1.3), y+Inches(0.1), Inches(13), Inches(0.4), "HR操作: 点击候选人 -> 展开查看原始证书图片 + 系统判定原因 -> 做出最终决定", WHITE, 13, True)

# HR decision pills -- BIG and prominent
hr_pill(s3, Inches(1.5), y+Inches(0.6), "PASS 通过", GREEN)
hr_pill(s3, Inches(5.5), y+Inches(0.6), "FAIL 淘汰", RED)
hr_pill(s3, Inches(9.5), y+Inches(0.6), "HOLD 待补充材料", ORANGE)

y += Inches(1.4)
big_header(s3, Inches(1), y, "数据看板", BLUE)
y += Inches(0.8)
box(s3, Inches(1), y, Inches(14), Inches(0.9))
txt(s3, Inches(1.3), y+Inches(0.08), Inches(13), Inches(0.3), "分类标签: 全部候选人 | 通过 | 不通过 | 待定 | 待核验", WHITE, 13, True)
txt(s3, Inches(1.3), y+Inches(0.45), Inches(13), Inches(0.3), "内联下拉框修改状态 -> 自动归类 | CSV导出 | 已判定可重新修改 | 招聘类型(校招/社招/实习)标签", GRAY, 11)

# ═══════════════════ SLIDE 4: 判定标准明细 ═══════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s4, BG)
txt(s4, Inches(0), Inches(0.08), Inches(16), Inches(0.5), "各核验环节判定标准明细", WHITE, 26, True, PP_ALIGN.CENTER)

def rule_row(s4, y, name, p, r, a):
    s = box(s4, Inches(0.8), y, Inches(2.3), Inches(0.62))
    tf = s.text_frame; tf.margin_left = Inches(0.1)
    p1 = tf.paragraphs[0]; p1.text = name; p1.font.size = Pt(10); p1.font.color.rgb = WHITE; p1.font.bold = True

    s = box(s4, Inches(3.3), y, Inches(3.9), Inches(0.62), border=GREEN)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.08)
    p1 = tf.paragraphs[0]; p1.text = p; p1.font.size = Pt(9); p1.font.color.rgb = GREEN

    s = box(s4, Inches(7.4), y, Inches(3.9), Inches(0.62), border=ORANGE)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.08)
    p1 = tf.paragraphs[0]; p1.text = r; p1.font.size = Pt(9); p1.font.color.rgb = ORANGE

    s = box(s4, Inches(11.5), y, Inches(3.7), Inches(0.62), border=RED)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.08)
    p1 = tf.paragraphs[0]; p1.text = a; p1.font.size = Pt(9); p1.font.color.rgb = RED

# Header
yh = Inches(0.85)
for i, (title, color) in enumerate([("检查项", BLUE), ("PASS 自动通过", GREEN), ("REVIEW 需复核", ORANGE), ("ALERT 高风险告警", RED)]):
    cx = Inches(0.8) + i * Inches(3.8) if i < 2 else Inches(7.4) + (i-2) * Inches(4.1)
    w = Inches(2.3) if i == 0 else Inches(3.9) if i < 3 else Inches(3.7)
    s = box(s4, cx, yh, w, Inches(0.45), color, color)
    tf = s.text_frame; p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

rules = [
    ("教育部名单", "学校在名单中", "--", "学校不在名单中，直接拒绝"),
    ("学信库查询", "查到记录且字段一致\n(姓名/学校/专业/学历/日期)", "查到记录但部分字段不一致", "未查到该学籍记录"),
    ("证书状态", "毕业证书 / 学位证书", "结业证书 / 肄业证书", "--"),
    ("毕业年龄", "年龄在正常范围\n本科18-30 硕士22-40 博士24-45", "年龄偏大或偏小", "--"),
    ("院校-专业一致性", "专业在学校公开列表中", "专业未在列表中\n(建议人工核实)", "--"),
    ("ELA图片篡改", "未检测到篡改痕迹\n(风险分 < 0.30)", "轻微异常(风险分 0.30-0.40)\n可能为翻拍/多次转发", "明显篡改痕迹(风险分 > 0.40)\n疑似P图修改"),
    ("双证交叉比对", "学位证与毕业证:\n姓名/学校/专业/学历/日期均一致", "任一项不一致", "--"),
    ("简历-学校匹配", "证书学校与简历学校一致\n(含曾用名/简称归一化)", "不一致", "--"),
    ("简历-专业匹配", "专业包含匹配通过\n(例: 法学 包含于 民法学)", "不匹配", "--"),
    ("简历-年份匹配", "证书毕业年份 = 简历结束年份\n(严格相等)", "年份不一致\n(可能延期/提前毕业)", "--"),
    ("简历缺证-校招", "--", "应届待取证: 标记不阻断", "--"),
    ("简历缺证-社招/实习", "--", "--", "已毕业学历必须有证\n缺证即不通过"),
    ("简历特殊状态", "ongoing/dropout/exchange/minor\n自动跳过，不要求证书", "--", "--"),
    ("综合判定", "所有检查均为PASS", "存在REVIEW且无ALERT", "存在至少一项ALERT"),
]

y = Inches(1.4)
for name, p, r, a in rules:
    rule_row(s4, y, name, p, r, a)
    y += Inches(0.68)

txt(s4, Inches(0.8), y+Inches(0.1), Inches(14), Inches(0.3), "判定优先级: ALERT > REVIEW > PASS | HR在人工核验台查看具体未通过原因后，做出最终判定(通过/淘汰/待补充材料)", GRAY, 10)

prs.save("D:/edu-verify/学历核验系统流程图_v4.pptx")
print("Saved -- 4 slides")
