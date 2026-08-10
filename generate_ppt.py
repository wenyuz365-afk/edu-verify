"""AI学历核验系统 - 汇报PPT vFinal 浅蓝科技感 15页"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ─── 浅蓝科技配色 ───
DARK  = RGBColor(0x0F, 0x2B, 0x45)
BLUE  = RGBColor(0x2D, 0x6F, 0xAB)
SKY   = RGBColor(0x5B, 0xA0, 0xCF)
LITE  = RGBColor(0xE8, 0xF2, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x3A, 0x3A, 0x3A)
MUTE  = RGBColor(0x8A, 0x9A, 0xA8)
WARN  = RGBColor(0xC4, 0x7A, 0x3A)
PASS  = RGBColor(0x3A, 0x8A, 0x6A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width; H = prs.slide_height
TOTAL = 18

def bg(s, c=WHITE):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def title_bar(s, text, y=Inches(0.5)):
    tb = s.shapes.add_textbox(Inches(0.8), y, W-Inches(1.6), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(24); p.font.color.rgb = BLUE; p.font.bold = True

def sub_bar(s, text, y=Inches(1.15)):
    tb = s.shapes.add_textbox(Inches(0.8), y, W-Inches(1.6), Inches(0.35))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(12); p.font.color.rgb = MUTE

def tx(s, text, l, t, w, h, sz=12, c=GRAY, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = bold; p.alignment = align
    p.line_spacing = Pt(int(sz*1.5))

def mtx(s, lines, l, t, w, h, sz=11, c=GRAY, sp=Pt(18)):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz); p.font.color.rgb = c; p.line_spacing = sp

def hline(s, l, t, w, c=BLUE, pw=2.5):
    ln = s.shapes.add_connector(1, l, t, l+w, t)
    ln.line.color.rgb = c; ln.line.width = Pt(pw)

def foot(s, n):
    tx(s, f"{n}/{TOTAL}", Inches(12.3), H-Inches(0.4), Inches(0.8), Inches(0.3), sz=8, c=MUTE)

def tbl(s, l, t, w, h, hdrs, rows, cw=None):
    nr = len(rows)+1; nc = len(hdrs)
    t_ = s.shapes.add_table(nr, nc, l, t, w, h).table
    if cw:
        for i, w_ in enumerate(cw): t_.columns[i].width = w_
    for j, hdr in enumerate(hdrs):
        cl = t_.cell(0, j); cl.text = hdr
        for p in cl.text_frame.paragraphs:
            p.font.size = Pt(10); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        cl.fill.solid(); cl.fill.fore_color.rgb = BLUE; cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cl = t_.cell(i+1, j); cl.text = str(val)
            for p in cl.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.CENTER
            cl.fill.solid(); cl.fill.fore_color.rgb = WHITE if i%2==0 else LITE
            cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i in range(nr):
        for j in range(nc):
            tcPr = t_.cell(i, j)._tc.get_or_add_tcPr()
            for bn in ['a:lnL','a:lnR','a:lnT','a:lnB']:
                tag = tcPr.find(qn(bn))
                if tag is None: tag = etree.SubElement(tcPr, qn(bn))
                tag.set('w','6350'); tag.set('cap','flat')
                sf = tag.find(qn('a:solidFill'))
                if sf is None: sf = etree.SubElement(tag, qn('a:solidFill'))
                sr = sf.find(qn('a:srgbClr'))
                if sr is None: sr = etree.SubElement(sf, qn('a:srgbClr'))
                sr.set('val','D0D8E0')
    return t_

def card(s, l, t, w, h, title, lines):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = LITE; sh.line.width = Pt(1.5)
    # 顶部色条
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tx(s, title, l+Inches(0.25), t+Inches(0.2), w-Inches(0.5), Inches(0.3), sz=14, c=BLUE, bold=True)
    mtx(s, lines, l+Inches(0.25), t+Inches(0.65), w-Inches(0.5), h-Inches(0.85), sz=10, c=GRAY)

def flow_node(s, l, t, w, h, text, fill=BLUE, tc=WHITE, sz=10):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    p = sh.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = tc; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

def arrow_r(s, l, t, w=Inches(0.3), h=Inches(0.18)):
    ss = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    ss.fill.solid(); ss.fill.fore_color.rgb = SKY; ss.line.fill.background()

def sec_page(s, num, title, st):
    bg(s, DARK)
    hline(s, Inches(1.5), Inches(2.8), Inches(1.5), SKY, 2)
    tx(s, num, Inches(1.5), Inches(2.0), Inches(10), Inches(0.5), sz=12, c=SKY, bold=True)
    tx(s, title, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8), sz=32, c=WHITE, bold=True)
    tx(s, st, Inches(1.5), Inches(3.9), Inches(10), Inches(0.5), sz=14, c=MUTE)

# ═══ 1. 封面 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
hline(s, Inches(1.5), Inches(2.6), Inches(1.5), SKY, 3)
tx(s, "学历信息 AI 自动核验系统", Inches(1.5), Inches(2.9), Inches(10), Inches(0.9), sz=36, c=WHITE, bold=True)
tx(s, "AI-Driven Credential Verification System", Inches(1.5), Inches(3.8), Inches(10), Inches(0.5), sz=14, c=SKY)
tx(s, "OCR 识别 · LLM 字段提取 · 三层递进核验 · 异常预警", Inches(1.5), Inches(4.6), Inches(10), Inches(0.4), sz=11, c=MUTE)
tx(s, "2026 年 7 月", Inches(1.5), Inches(6.0), Inches(3), Inches(0.3), sz=10, c=MUTE)
foot(s, 1)

# ═══ 2. 目录 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "目录"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)
items = [
    ("01", "项目背景与目标", "校招学历核验现状 · AI 解决思路 · 项目定位"),
    ("02", "后端技术架构", "三层递进核验引擎 · 预处理管线 · 综合判定逻辑"),
    ("03", "技术栈", "OCR · LLM · 数据库 · 图像处理 · 部署方案"),
    ("04", "前端交互设计", "上传核验 · 人工核验台 · 数据看板 · 工作流"),
    ("05", "关键决策与特殊情况", "end_year 判定 · 去重策略 · 阈值校准 · 退学/交换/辅修"),
    ("06", "Demo 演示", "测试场景 · 核验路径 · 预期结果"),
]
for i, (num, tt, desc) in enumerate(items):
    y = Inches(1.7) + i*Inches(0.85)
    tx(s, num, Inches(1.2), y, Inches(0.5), Inches(0.3), sz=16, c=BLUE, bold=True)
    tx(s, tt, Inches(1.9), y, Inches(4), Inches(0.3), sz=14, c=BLUE, bold=True)
    tx(s, desc, Inches(6.2), y+Inches(0.02), Inches(6), Inches(0.3), sz=10, c=GRAY)
foot(s, 2)

# ═══ 3. Part1 分页 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); sec_page(s, "Part 1", "项目背景与目标", "校招学历核验的现状与 AI 解决方案"); foot(s, 3)

# ═══ 4. 背景与痛点 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "校招学历核验 — 现状与挑战"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)
mtx(s, ["校园招聘季，HR 需在短时间内处理数百份候选人材料。学历信息真伪核验是不可或缺的一环，但传统人工方式面临多重困境。"],
    Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7), sz=10, c=GRAY)
tbl(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(3.5),
    ["痛点", "具体表现", "后果"],
    [["证书格式多样", "学位证/毕业证/国外学历/军校/成教，排版各异", "人工逐行读取效率极低"],
     ["批量处理压力", "校招季数百份证书需逐一核对", "单份耗时 10-15 分钟"],
     ["信息孤岛", "简历、学位证、毕业证信息独立", "造假者可利用信息不对称"],
     ["伪造手段升级", "PS 篡改、虚假院校、冒名顶替", "肉眼无法识别图片级篡改"],
     ["缺乏统一工具", "现有方案多系统切换，无统一核验平台", "流程碎片化，易遗漏"],
     ["院校知识盲区", "HR 难掌握全国院校及专业情况", "虚假院校/专业难发现"]],
    [Inches(1.8), Inches(5.0), Inches(4.7)])

tx(s, "项目目标", Inches(0.8), Inches(6.2), Inches(3), Inches(0.3), sz=13, c=BLUE, bold=True)
mtx(s, ["构建 AI 驱动的学历自动验证系统，覆盖证书识别 → 字段提取 → 多维核验 → 异常预警全流程",
    "将 HR 单人核验时间压缩至 30 秒以内，效率提升 20 倍以上"],
    Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.7), sz=10, c=GRAY)
foot(s, 4)

# ═══ 5. Part2 分页 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); sec_page(s, "Part 2", "后端技术架构", "三层递进核验引擎 · 预处理管线 · 综合判定"); foot(s, 5)

# ═══ 6. 架构全景 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "系统架构全景"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)

# 流程横向
steps = [("上传证书", BLUE), ("OCR 识别", BLUE), ("LLM 提取", BLUE), ("L1 篡改检测", DARK), ("L2 学信库", DARK), ("L3 交叉核验", DARK), ("判定输出", BLUE)]
nx = Inches(0.3); nw = Inches(1.6); nh = Inches(0.8); ny = Inches(1.7); ng = Inches(0.2)
for i, (label, clr) in enumerate(steps):
    x = nx + i*(nw+ng)
    flow_node(s, x, ny, nw, nh, label, fill=clr, sz=9)
    if i < len(steps)-1:
        arrow_r(s, x+nw+Inches(0.02), ny+nh/2-Inches(0.09), Inches(0.16))

# 三层详解
tx(s, "三层递进核验引擎（全部执行后综合判定）", Inches(0.8), Inches(2.9), Inches(6), Inches(0.3), sz=12, c=BLUE, bold=True)
layers = [
    ("L1: 图片篡改检测", "ELA 分析所有证书图片。high → 异常(高) / medium → 异常(中) / low → 通过", SKY),
    ("L2: 学信库核验", "六字段逐一匹配：姓名/学校/专业/学历层级/毕业日期/证书编号。不匹配 → 异常(高) / 未查到 → 异常(中) / 全匹配 → 通过", BLUE),
    ("L3: 交叉核验", "双证比对 + 简历验证 + 毕业状态判定(end_year)。退学/交换/辅修 → 直接跳过", DARK),
]
for i, (tt, desc, clr) in enumerate(layers):
    y = Inches(3.4) + i*Inches(1.2)
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.1), Inches(1.0))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = LITE; rect.line.width = Pt(1)
    tx(s, tt, Inches(1.0), y+Inches(0.1), Inches(4), Inches(0.3), sz=12, c=clr, bold=True)
    tx(s, desc, Inches(1.0), y+Inches(0.5), Inches(11), Inches(0.4), sz=10, c=GRAY)

tx(s, "综合判定: 三层全部执行 → 任一高严重度 → REVIEW 高危 / 任一中低异常 → REVIEW 需复核 / 全部通过 → PASS",
    Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.3), sz=10, c=BLUE, bold=True)
foot(s, 6)

# ═══ 7. OCR + LLM ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "预处理管线 — OCR & LLM 字段提取"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)

# OCR 流程
tx(s, "OCR 识别管线", Inches(0.8), Inches(1.7), Inches(4), Inches(0.3), sz=12, c=BLUE, bold=True)
ocr_steps = ["原始输入\n(JPG/PNG/PDF)", "PDF→PNG", "HSV 印章分离", "自适应缩放", "PaddleOCR\n文字识别"]
ox = Inches(0.3); ow = Inches(2.2); oh = Inches(0.8); oy = Inches(2.2)
for i, label in enumerate(ocr_steps):
    clr = BLUE if i not in [2,3] else DARK
    flow_node(s, ox + i*(ow+Inches(0.25)), oy, ow, oh, label, fill=clr, sz=9)
    if i < 4: arrow_r(s, ox + (i+1)*(ow+Inches(0.25))-Inches(0.27), oy+oh/2-Inches(0.09), Inches(0.2))

# LLM 对比
tx(s, "LLM 字段提取 — 纯大模型驱动", Inches(0.8), Inches(3.5), Inches(5), Inches(0.3), sz=12, c=BLUE, bold=True)
tbl(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(2.0),
    ["对比维度", "传统正则方案", "LLM 方案（本项目）"],
    [["通用性", "每种证书格式写规则", "自然语言理解，零正则，任意排版"],
     ["中文年份识别", "无法识别\"二〇二五\"", "Prompt 内置对照表，LLM 转数字"],
     ["键名兼容", "仅支持英文键名", "中英文自动映射（KEY_MAP）"],
     ["维护成本", "新格式需追加正则", "调优 Prompt，一次搞定"]],
    [Inches(2.0), Inches(4.8), Inches(4.7)])

tx(s, "提取字段: 姓名 | 学校 | 专业 | 学历层级 | 毕业日期 | 证书编号", Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.3), sz=10, c=BLUE, bold=True)
foot(s, 7)

# ═══ 8. Layer 1 详解 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "Layer 1 — 图片篡改检测（ELA）"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)
mtx(s, ["第一层核验，对所有上传证书图片进行 Error Level Analysis 错误级别分析。"],
    Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5), sz=10, c=GRAY)
tbl(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(3.5),
    ["项目", "说明"],
    [["原理", "JPEG 图片经 PS 修改后，不同区域的压缩错误率不一致。通过分析各区域高频分量差异，识别被篡改区域。"],
     ["HSV 印章预处理", "证书上的红色印章本身会产生高频分量导致误报。先用 HSV 色彩空间分离并移除印章区域，再对剩余内容做 ELA 分析。"],
     ["阈值经实测校准", "真证 0.45 → low · 假证 0.66 → high。阈值设于 high=0.52 / medium=0.46，精准区分。"],
     ["high 风险", "疑似高危篡改 → 标记异常(高)，继续后续核验"],
     ["medium 风险", "图片存在轻微异常 → 标记异常(中)，继续后续核验"],
     ["low 风险", "未检测到明显篡改痕迹 → 通过，继续后续核验"]],
    [Inches(2.5), Inches(9.0)])
tx(s, "无论结果如何，Layer 1 执行完毕后继续 Layer 2，不提前终止。最终三层汇总后综合判定。",
    Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.3), sz=10, c=BLUE, bold=True)
foot(s, 8)

# ═══ 9. Layer 2 详解 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "Layer 2 — 学信库核验（六字段匹配）"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)
mtx(s, ["第二层核验，每份证书的 LLM 提取字段与学信库记录逐一比对。"],
    Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5), sz=10, c=GRAY)

tx(s, "六字段比对", Inches(0.8), Inches(2.3), Inches(4), Inches(0.3), sz=12, c=BLUE, bold=True)
tbl(s, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3.2),
    ["字段", "比对方式"],
    [["姓名", "精确匹配，支持曾用名模糊匹配"],
     ["学校", "精确匹配，支持院校曾用名（如\"北京钢铁学院\"→\"北京科技大学\"）"],
     ["专业", "精确匹配，支持包含匹配"],
     ["学历层级", "匹配本科/硕士/博士，支持等效映射（学士→本科）"],
     ["毕业日期", "YYYY-MM 标准化后精确匹配"],
     ["证书编号", "OCR 提取编号与学信库记录比对"]],
    [Inches(1.3), Inches(4.2)])

tx(s, "判定规则", Inches(7.0), Inches(2.3), Inches(4), Inches(0.3), sz=12, c=BLUE, bold=True)
tbl(s, Inches(7.0), Inches(2.7), Inches(5.5), Inches(3.2),
    ["结果", "处理"],
    [["六字段全部一致", "通过，继续 Layer 3"],
     ["任一字段不匹配", "标记异常(高)，继续 Layer 3"],
     ["学信库未查到此人", "标记异常(中)，继续 Layer 3"]],
    [Inches(2.2), Inches(3.3)])

tx(s, "Layer 2 不提前终止。无论匹配结果如何，继续执行 Layer 3。",
    Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.3), sz=10, c=BLUE, bold=True)
foot(s, 18)

# ═══ 10. Layer 3 详解 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "Layer 3 — 交叉核验"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)
mtx(s, ["第三层核验，通过前两层的证书在此进行多维度综合交叉验证。"],
    Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5), sz=10, c=GRAY)

tbl(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(4.5),
    ["检查项", "方法", "异常条件", "严重度"],
    [["双证交叉比对", "学位证 vs 毕业证：姓名/学校/专业/学历层级/毕业日期\n逐字段比对", "任一不匹配", "高"],
     ["简历交叉验证", "简历声明 vs 证书提取：学校/专业/年份比对", "不一致", "中"],
     ["毕业状态判定", "end_year < 今年：已毕业缺证\nend_year == 今年：应届待取证\nend_year > 今年：在读正常\nend_year 缺失：无法判定", "已毕业缺证\n应届缺证\n在读\n无法判定", "高\n低\nPASS\n中"],
     ["退学/交换/辅修", "LLM 提取 status 字段命中则跳过", "—", "PASS"],
     ["简历-证书交叉比对", "简历有匹配证书时：比对学校/专业/年份一致性", "不一致", "中"]],
    [Inches(1.8), Inches(5.2), Inches(3.3), Inches(1.2)])

tx(s, "综合判定: 三层全部检查项汇总 → 任一高严重度 → REVIEW 高危 / 任一中低异常 → REVIEW 需复核 / 全部 PASS → 核验通过",
    Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.3), sz=10, c=BLUE, bold=True)
foot(s, 16)

# ═══ 11. Part3 分页 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); sec_page(s, "Part 3", "技术栈", "OCR · LLM · 数据库 · 图像处理 · 部署"); foot(s, 17)

# ═══ 9. 技术栈 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "技术栈全览"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)

tbl(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(4.8),
    ["层级", "技术", "版本/参数", "选型原因"],
    [["OCR引擎", "PaddleOCR", "2.9.1 (PaddlePaddle 2.6.2)", "本地离线识别，零 API 配额，印章去除专项优化"],
     ["AI模型", "DeepSeek V4 Pro", "thinking=disabled, max_tokens=4096", "Anthropic 兼容接口，纯 LLM 字段提取，零正则规则"],
     ["UI框架", "Streamlit", "1.59", "Python 全栈，一个月交付周期，三步核验交互流"],
     ["数据库", "SQLite", "edu_verify.db (2.4MB)", "单机 Demo，零配置，证书编号唯一去重"],
     ["图像处理", "OpenCV", "HSV 色彩空间", "印章分离、自适应缩放、图像增强"],
     ["篡改检测", "自研 ELA", "high=0.52 / med=0.46", "Error Level Analysis，经测试集校准阈值"],
     ["PDF处理", "pypdfium2", "scale=2", "PDF 渲染为 PNG 供 OCR 处理"],
     ["部署", "Docker + ngrok", "Python 3.12.9", "容器化部署，内网穿透公网演示"]],
    [Inches(1.5), Inches(2.5), Inches(3.5), Inches(4.0)])
foot(s, 18)

# ═══ 10. Part4 分页 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); sec_page(s, "Part 4", "前端交互设计", "上传核验 · 人工核验台 · 数据看板"); foot(s, 16)

# ═══ 11. 上传核验 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "上传核验 — 两步交互流"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)

card(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.5),
    "Step 1: 材料提交",
    ["学位证 + 毕业证 + 简历（均必传）",
     "多学历层次: 本科/硕士/博士 灵活增减（最多 3 层）",
     "上传模式: 单个候选人 或 ZIP 批量导入",
     "AI 自动判定毕业状态，无需 HR 手动选择"])

card(s, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.5),
    "Step 2: 结果展示",
    ["判定总览: PASS/REVIEW 徽章 + 综合结论",
     "三层核验详情: 每层折叠展开，全部检查项可见",
     "通过项和异常项均展示，附带计数统计",
     "证书字段摘要表 + 简历交叉验证"])

card(s, Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.0),
    "批量核验",
    ["ZIP 压缩包上传，自动按文件夹分组",
     "按文件名关键词分类: 学位/毕业/简历",
     "批量结果支持 CSV 导出"])
card(s, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.0),
    "证书编号去重",
    ["同编号再次核验 → 自动覆盖上一条",
     "数据看板一人一行，不重复"])
foot(s, 17)

# ═══ 12. 人工核验台 & 数据看板 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "人工核验台 & 数据看板"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)

card(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.8),
    "人工核验台",
    ["待核验列表: 自动筛选未审核候选人，按严重度排列",
     "问题项明细: 仅显示异常检查项，逐条列出原因",
     "材料查看: 学位证/毕业证/简历原图可展开",
     "HR 判定: 通过/不通过/待定，每次判定需备注",
     "支持回溯修改 + 补充上传材料"])

card(s, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.8),
    "数据看板",
    ["四大指标卡片: 总人数/通过/不通过/待定",
     "5 个 Tab: 全部/通过/不通过/待定/待核验",
     "在线编辑: data_editor 直接修改判定",
     "按分类 CSV 导出，文件名自动带日期戳",
     "一人一行: 证书编号去重保证不重复"])

# 工作流
tx(s, "HR 完整工作流", Inches(0.8), Inches(4.9), Inches(3), Inches(0.3), sz=12, c=BLUE, bold=True)
wf = ["批量上传", "AI 自动核验\n(<30秒/人)", "分级: PASS\n或 REVIEW", "人工复核\n异常项", "导出名单"]
wx = Inches(0.5); ww = Inches(2.2); wh = Inches(0.7); wy = Inches(5.5)
for i, label in enumerate(wf):
    clr = BLUE if i in [0,3,4] else DARK
    flow_node(s, wx+i*(ww+Inches(0.25)), wy, ww, wh, label, fill=clr, sz=9)
    if i < 4: arrow_r(s, wx+(i+1)*(ww+Inches(0.25))-Inches(0.27), wy+wh/2-Inches(0.09), Inches(0.2))
foot(s, 18)

# ═══ 13. Part5 分页 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); sec_page(s, "Part 5", "关键决策与特殊情况", "end_year 判定 · 去重 · 阈值 · 边界处理"); foot(s, 16)

# ═══ 14. 关键决策 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, "关键决策与特殊情况处理"); hline(s, Inches(0.8), Inches(1.15), Inches(2), SKY, 2)

tbl(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.0),
    ["决策点", "方案", "理由"],
    [["毕业状态判定", "end_year < 今年→已毕业 / == 今年→应届 / > 今年→在读\n退学/交换/辅修: 直接 PASS，不参与判定", "替代人工选择招聘类型，减少 HR 操作"],
     ["缺证判定", "已毕业缺证→异常(高) / 应届缺证→异常(低) / 在读缺证→通过", "毕业年份是唯一客观标准"],
     ["证书编号去重", "同编号核验 → UPDATE 覆盖旧记录\n数据库一人一行", "证书编号一对一，天然的 unique key"],
     ["ELA 阈值校准", "high=0.52, med=0.46，经真证/假证实测校准\n假证 0.66→high / 真证 0.45→low", "避免扫描件印章产生误报"],
     ["三层全跑策略", "L1/L2/L3 全部执行完毕后综合判定\n不因单层异常提前终止", "HR 需要完整画像，不是分段信息"],
     ["简历无法提取年份", "end_year 缺失或非数字 → REVIEW\n提示 HR 手动确认", "不能硬猜，宁可转到人工"],
     ["退学/交换/辅修", "LLM 提取 status 字段，命中即跳过核验\n不需要证书", "这些情况本来就不该有证书"],
     ["两档判定体系", "PASS + REVIEW(高/中/低)\n替代原 PASS/REVIEW/ALERT 三档", "REVIEW/ALERT 对 HR 动作相同，合并为严重度区分"]],
    [Inches(2.2), Inches(5.5), Inches(4.6)])
foot(s, 17)

# ═══ 15. Demo + 总结 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
hline(s, Inches(1.5), Inches(0.8), Inches(1.5), SKY, 2)
tx(s, "Demo 演示 & 总结", Inches(1.5), Inches(1.0), Inches(10), Inches(0.7), sz=28, c=WHITE, bold=True)

tx(s, "演示场景", Inches(1.5), Inches(2.0), Inches(5), Inches(0.4), sz=14, c=SKY, bold=True)
mtx(s, [
    "场景 1: 正常通过 — 真证 + 匹配简历 → 三层全 PASS",
    "场景 2: 图片篡改 — 假证 → L1 异常(高) + L2/L3 同时检测 → 综合 REVIEW",
    "场景 3: 学信库不匹配 — 真证但非学信库人员 → L2 异常 + L3 验证 → REVIEW",
    "场景 4: 批量上传 — ZIP 包含多人 → 自动分组 → 逐人三层核验",
], Inches(1.5), Inches(2.5), Inches(5.5), Inches(2.0), sz=10, c=WHITE)

tx(s, "技术成果", Inches(7.5), Inches(2.0), Inches(5), Inches(0.4), sz=14, c=SKY, bold=True)
mtx(s, [
    "完整的 OCR → LLM → 三层核验 → 判定输出全流程闭环",
    "583 国内 + 191 海外院校白名单",
    "end_year 自动判定毕业状态，零人工输入",
    "证书编号唯一去重，自动覆盖",
    "20+ 轮迭代，环境/OCR/AI/数据/逻辑/UI 全覆盖",
], Inches(7.5), Inches(2.5), Inches(5.5), Inches(2.0), sz=10, c=WHITE)

tx(s, "后续方向", Inches(1.5), Inches(4.8), Inches(5), Inches(0.4), sz=14, c=SKY, bold=True)
mtx(s, [
    "接入学信网官方 API，实现权威数据源实时查询",
    "接入教育部留学服务中心认证接口",
    "企业 SSO 集成，嵌入招聘审批流程",
], Inches(1.5), Inches(5.3), Inches(5.5), Inches(1.5), sz=10, c=WHITE)

hline(s, Inches(1.5), Inches(6.3), Inches(10), MUTE, 1)
tx(s, "Python 3.12 · PaddleOCR 2.6 · DeepSeek V4 Pro · Streamlit 1.59 · SQLite · Docker",
    Inches(1.5), Inches(6.5), Inches(11), Inches(0.3), sz=9, c=MUTE)
foot(s, 18)

# ─── 保存 ───
out = "D:/edu-verify/AI学历核验系统_汇报PPT.pptx"
prs.save(out)
print(f"OK: {out}")
